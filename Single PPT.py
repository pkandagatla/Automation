#!/usr/bin/env python
# coding: utf-8

# In[10]:


from pptx import Presentation  
from pptx.chart.data import CategoryChartData  
from pptx.table import Table, _Row, _Column, _Cell
from pptx.enum.chart import XL_CHART_TYPE  
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pandas as pd
import numpy as np
from openpyxl import load_workbook
from pptx.dml.color import RGBColor
from datetime import datetime
import win32com.client
import warnings
warnings.filterwarnings("ignore")


def coordinate_finder_insheet_df(df,valuetofind,recurrence):   
    found=False
    count=0
    for i in range(len(df)):
        for j in range(len(df.columns)):
            if(df.iloc[i,j]==valuetofind):
                row_no=i+1
                col_no=j+1
                count+=1
                if(count==recurrence):
                    found=True
                    break
        if found is True:
            break 
    return(row_no,col_no)

def find_replace_text(presentation, find_text, replace_text): 
    for slide in presentation.slides: 
        for shape in slide.shapes: 
            if shape.has_text_frame: 
                for paragraph in shape.text_frame.paragraphs: 
                    for run in paragraph.runs: 
                        if find_text in run.text: 
                            run.text = run.text.replace(find_text, replace_text) 

def highlight_format(shape):      
    tf = shape.text_frame
    p = tf.paragraphs[0]
    font = p.font
    font.name = 'Arial'
    font.bold = True
    p.alignment = PP_ALIGN.CENTER

def text_format_no_align(shape):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    font = p.font
    font.name = 'Arial'
    font.size = Pt(10.5)

def text_format_heading(shape):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    font = p.font
    font.name = 'Arial'
    font.size = Pt(10.5)
    font.bold = True
    p.alignment = PP_ALIGN.CENTER
    font.color.rgb = RGBColor(255,255,255)

def remove_row(table: Table,
               row_to_delete: _Row):
    table._tbl.remove(row_to_delete._tr)
    
def bold(shape):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    font = p.font
    font.bold = True
    
def text_format(shape):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    font = p.font
    font.name = 'Arial'
    font.size = Pt(10.5)
    p.alignment = PP_ALIGN.CENTER
    
def text_format_footer(shape):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    font = p.font
    font.name = 'Arial'
    font.size = Pt(9)
    font.color.rgb = RGBColor(255,255,255)
    
def font_color(shape,color):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    font = p.font
    if color == "green":
        font.color.rgb = RGBColor(0,176,80)
    else:
        font.color.rgb = RGBColor(255,0,0)
        
def mi(mi):
    if mi > 0: idm = "below"
    else: idm = "above"
    return idm

def qhcr(qhcr):
    if qhcr > 0: idq = "below"
    else: idq = "above"
    return idq

def qhcr_color(qhcr):
    if qhcr > 0: color = "red"
    else: color = "green"
    return color

def qhcr_color_opp(qhcr):
    if qhcr > 0: color = "green"
    else: color = "red"
    return color

template_path = r"C:\Users\pragna_kandagatla\Desktop\Automation\Assessment PPT\AR Assessment - PPT - template.pptm"
excel_path = r"C:\Users\pragna_kandagatla\Desktop\Automation\Assessment PPT\test.xlsx"
output_path = r"C:\Users\pragna_kandagatla\Downloads\powerpoint.pptm"
final_path = r"C:\Users\pragna_kandagatla\Downloads\AR Assessment - PPT.pptx"
root = Presentation(template_path) 

assessment_name=input("Please enter the title of the assessment = ")
template_assessment_metrics=pd.read_excel(excel_path,sheet_name="Assessment Metrics",header=None)
template_aging_cal=pd.read_excel(excel_path,sheet_name="Aging Cal",header=None)
template_roi =pd.read_excel(excel_path,sheet_name="ROI Calculator",header=None)

template = pd.ExcelFile(excel_path)
if "DSO & NCR" in template.sheet_names:
    template_ncr = pd.read_excel(excel_path,sheet_name="DSO & NCR",header=None)
else:
    template_ncr = pd.read_excel(excel_path,sheet_name="NCR",header=None)
    
#----- replacing assessment name
find_replace_text(root, "Assessment Name", assessment_name) 


#--------------------------------------------------------------------Slide 5

slide = root.slides[4]

#---------NCR-------------
row_no, col_no = coordinate_finder_insheet_df(template_assessment_metrics,"Net Collection Rate",1)
ncr = template_assessment_metrics.iloc[row_no-1,col_no]

row_no, col_no = coordinate_finder_insheet_df(template_assessment_metrics,"Net Collection Rate excluding Private",1)
ncr_wo_private = template_assessment_metrics.iloc[row_no-1,col_no]

row_no, col_no = coordinate_finder_insheet_df(template_assessment_metrics,"Industry Standard",1)
missed_mi = template_assessment_metrics.iloc[row_no-1,col_no+1]

row_no, col_no = coordinate_finder_insheet_df(template_assessment_metrics,"QHCR Benchmark",1)
missed_qhcr = template_assessment_metrics.iloc[row_no-1,col_no+1]

ncr_mi = 0.95 - ncr
ncr_qhcr =0.98 - ncr

idq = qhcr(ncr_qhcr)
idm = mi(ncr_mi)
color = qhcr_color(ncr_qhcr)

if missed_mi == missed_qhcr == "-" : missed = "$0K"
elif missed_mi == "-": missed = "${:,.0f}K".format(missed_qhcr/1000)
elif missed_qhcr == "-": missed = "${:,.0f}K".format(missed_mi/1000)
else: missed = "${:,.0f}K".format(missed_mi/1000) + " - " + "${:,.0f}K".format(missed_qhcr/1000)

ncr_text = assessment_name + " Net Collection Percentage is " + "{:.1%}".format(abs(ncr_qhcr)) + " " + idq + " QHCR’s Benchmark and is "\
                     + "{:.1%}".format(abs(ncr_mi)) + " " + idm + " industry standard, representing " + missed + " in missed collections yearly"

for shape in slide.shapes:
    if shape.name == "TextBox 5":
        shape.text = "{:.1%}".format(abs(ncr_qhcr))
        highlight_format(shape)
        font_color(shape,color)
        
    if shape.name == "TextBox 3":
        shape.text = ncr_text
        text_format_no_align(shape)
        
    if shape.name == "Oval 33":
        line = shape.line
        if idq == "below": line.color.rgb = RGBColor(255,0,0)
        else: line.color.rgb = RGBColor(0,176,80) 
        
#--------------------DSO------------------
        
row_no, col_no = coordinate_finder_insheet_df(template_assessment_metrics,"DSO",1)
dso = round(template_assessment_metrics.iloc[row_no-1,col_no])

dso_mi = round(40 - dso)
dso_qhcr = round(34 - dso)

idq = qhcr(dso_qhcr)
idm = mi(dso_mi)
color = qhcr_color_opp(dso_qhcr)

dso_text = assessment_name + " Days of Sales Outstanding is " + str(abs(dso_qhcr)) + " days " + idq +\
                                  " QHCR’s Benchmark and is " + str(abs(dso_mi)) + " days "+ idm + " industry standard"

for shape in slide.shapes:
    if shape.name == "TextBox 44":
        shape.text = str(abs(dso_qhcr))
        highlight_format(shape)
        font_color(shape,color)
        
    if shape.name == "TextBox 4":
        shape.text = dso_text
        text_format_no_align(shape)
        
    if shape.name == "Oval 43":
        line = shape.line
        if idq == "below": line.color.rgb = RGBColor(0,176,80)
        else: line.color.rgb = RGBColor(255,0,0) 
        
#----------------AR Over 90--------------------
        
row_no, col_no = coordinate_finder_insheet_df(template_assessment_metrics,"% AR Over 90",1)
ar = template_assessment_metrics.iloc[row_no-1,col_no]

ar_mi = 0.15 - ar
ar_qhcr = 0.11 - ar

idq = qhcr(ar_qhcr)
idm = mi(ar_mi)
color = qhcr_color_opp(ar_qhcr)

ar_text = assessment_name +  " % of A/R Over 90 is " + "{:.1%}".format(abs(ar_qhcr)) + " " + idq +\
                                  " QHCR’s Benchmark and is " + "{:.1%}".format(abs(ar_mi)) + " " + idm + " industry standard"

for shape in slide.shapes:
    if shape.name == "TextBox 48":
        shape.text = "{:.1%}".format(abs(ar_qhcr))
        highlight_format(shape)
        font_color(shape,color)
        
    if shape.name == "TextBox 7":
        shape.text = ar_text
        text_format_no_align(shape)
        
    if shape.name == "Oval 47":
        line = shape.line
        if idq == "below": line.color.rgb = RGBColor(0,176,80)
        else: line.color.rgb = RGBColor(255,0,0) 
        
#----- date range

row_no, col_no = coordinate_finder_insheet_df(template_ncr,"Start Month",1)
smonth = template_ncr.iloc[row_no,col_no-1]

row_no, col_no = coordinate_finder_insheet_df(template_ncr,"End Month",1)
emonth = template_ncr.iloc[row_no,col_no-1]

for shape in slide.shapes:
    if shape.name == "TextBox 24":
        shape.text = "Report based on the information provided by " + assessment_name + " for the period " + smonth.strftime("%b'%y") + \
        " through " + emonth.strftime("%b'%y")+ "." +" " + "Net Collections % is a mean of " + smonth.strftime("%b'%y") + " to " + \
        emonth.strftime("%b'%y") + " to provide a more accurate calculation." + " "+ "Miscellaneous Applied Cash had been removed for the purpose of this analysis to provide a more accurate view of Net Collections."
        text_format_footer(shape)

for shape in slide.shapes:
    if shape.name == "Chart 10":
        chart = shape.chart
        
chart_data = CategoryChartData()
chart_data.categories = ["Net Collection %"]
chart_data.add_series(assessment_name + " - w/Private", [ncr])
chart_data.add_series(assessment_name + " - w/o Private", [ncr_wo_private])
chart_data.add_series("Industry Benchmark", [0.95])
chart_data.add_series("QHCR Benchmark", [0.98])
chart.replace_data(chart_data)
chart.has_3d = True

for shape in slide.shapes:
    if shape.name == "Chart 11":
        chart = shape.chart
        
chart_data = CategoryChartData()
chart_data.categories = ["DSO"]
chart_data.add_series(assessment_name, [dso])
chart_data.add_series("Industry Benchmark", [40])
chart_data.add_series("QHCR Benchmark" + " - w/o Private", [34])
chart.replace_data(chart_data)
chart.has_3d = True

for shape in slide.shapes:
    if shape.name == "Chart 12":
        chart = shape.chart
        
chart_data = CategoryChartData()
chart_data.categories = ["% AR over 90"]
chart_data.add_series(assessment_name, [ar])
chart_data.add_series("Industry Benchmark", [0.15])
chart_data.add_series("QHCR Benchmark", [0.11])
chart.replace_data(chart_data)
chart.has_3d = True

#--------------------------------------------------------------------Slide 6

slide = root.slides[5]

for shape in slide.shapes:
    if shape.name == "Chart 3":
        chart = shape.chart
        
chart_data = CategoryChartData()
chart_data.categories = [assessment_name]
chart_data.add_series("w/Private", [ncr])
chart_data.add_series("w/o Private", [ncr_wo_private])
chart.replace_data(chart_data)


#--------------------------------------------------------------------Slide 7

slide = root.slides[6]

row_no1, col_no1 = coordinate_finder_insheet_df(template_ncr,"Net Collection Rate",1)
row_no2, col_no2 = coordinate_finder_insheet_df(template_ncr,"Expected Revenue - Overall",1)
row_no3, col_no3 = coordinate_finder_insheet_df(template_ncr,"Payments - Overall",1)
row_no4, col_no4 = coordinate_finder_insheet_df(template_ncr,"Missed Collections - QHCR Benchmark",1)
row_no5, col_no5 = coordinate_finder_insheet_df(template_ncr,smonth,2)
row_no6, col_no6 = coordinate_finder_insheet_df(template_ncr,emonth,2)

ncr_values = list(template_ncr.iloc[row_no1-1,col_no5-1:col_no6])
charges_values = list(template_ncr.iloc[row_no2-1,col_no5-1:col_no6])
payments_values = list(template_ncr.iloc[row_no3-1,col_no5-1:col_no6])
month_values = list(template_ncr.iloc[row_no5-1,col_no5-1:col_no6])
missed_values = list(template_ncr.iloc[row_no4-1,col_no5-1:col_no6])
qhcr_values = [0.98]*len(ncr_values)

df = pd.DataFrame([month_values,ncr_values,charges_values,payments_values]).T
df.columns = ["Month","NCR","Charges","Payments"]

hmonth = df[df["NCR"] == df.max()["NCR"]].iloc[0,0]
lmonth = df[df["NCR"] == df.min()["NCR"]].iloc[0,0]

sncr = df["NCR"][0]
encr = df["NCR"][len(df["NCR"])-1]
sncrf = "{:.0f}%".format(sncr*100)
encrf = "{:.0f}%".format(encr*100)
cncr = "{:.0f}%".format(abs(sncr - encr)*100)

sc = df["Charges"][0]
ec = df["Charges"][len(df["NCR"])-1]
scf = "${:,.0f}K".format(sc/10**3)
ecf = "${:,.0f}K".format(ec/10**3)
changec = "${:,.0f}K".format(abs(sc - ec)/10**3)

sp = df["Payments"][0]
ep = df["Payments"][len(df["NCR"])-1]
changep = "${:,.0f}K".format(abs(sp - ep)/10**3)

if sncr - encr >= 0 :id1 = "decreased"
else: id1 = "increased"
    
if sc - ec >= 0 :id2 = "decreased"
else: id2 = "increased"    
    
if sp - ep >= 0 : id3 = "decreased"
else: id3 = "increased"   
    
for shape in slide.shapes:
    if shape.name == "Table 6":
        table = shape.table
    elif shape.name == "Chart 1":
        chart = shape.chart

for i in range(len(missed_values)):
    cell = table.cell(0,i+1)
    if type(missed_values[i]) != str :
        cell.text = "${:,.0f}".format(missed_values[i])
    else:
        cell.text = missed_values[i]
    text_format(cell)

chart_data = CategoryChartData()
chart_data.categories = month_values
chart_data.add_series("QHCR", qhcr_values)
chart_data.add_series(assessment_name, ncr_values)
chart.replace_data(chart_data)

main_bullet_points = [
    "Healthy net collection rate through " + hmonth.strftime("%b'%y") + "; lowest collection rate recorded in " + lmonth.strftime("%b'%y") +".",
    "Overall, net collection rate " + id1 + " from " + sncrf + " in " + smonth.strftime("%b'%y") + " to " + encrf +  " in " + emonth.strftime("%b'%y") + "."
]

sub_bullet_points = [ "While revenue " + id2 + " from " + scf + " in " + smonth.strftime("%b'%y") + " to " + ecf + " in " + emonth.strftime("%b'%y") + "(" + changec + ")" \
" during that period, payments " + id3 + " by " + changep + ".", "Missed revenue due to collection decline ~" + missed + "."]

for shape in slide.shapes:
    if shape.name == "Text Placeholder 54":
        for i, main_bullet in enumerate(main_bullet_points):
            if i == 0: p = shape.text_frame
            else : p = shape.text_frame.add_paragraph()
            p.text = main_bullet
            p.level = 0
        for sub_bullet in sub_bullet_points:
            p = shape.text_frame.add_paragraph()
            p.text = sub_bullet
            p.level = 1 
            font = p.font
            font.size = Pt(18)
            
#----------------------------------------------------Slide 8

slide = root.slides[7]

for shape in slide.shapes:
    if shape.name == "Chart 4":
        chart = shape.chart
        
chart_data = CategoryChartData()
chart_data.categories = ["A/R Over 90"]
chart_data.add_series(assessment_name, [ar])
chart_data.add_series("Minimum Industry Standard", [0.15])
chart_data.add_series("QHCR Benchmark", [0.11])
chart.replace_data(chart_data)


#--------------------------------------------------------------------Slide 10

slide = root.slides[9]

for shape in slide.shapes:
    if shape.name == "Chart 4":
        chart = shape.chart
        
chart_data = CategoryChartData()
chart_data.categories = ["DSO"]
chart_data.add_series(assessment_name, [dso])
chart_data.add_series("Minimum Industry Standard", [40])
chart_data.add_series("QHCR Benchmark", [34])
chart.replace_data(chart_data)


#--------------------------------------------------------------------Slide 9

slide = root.slides[8]
row_no, col_no = coordinate_finder_insheet_df(template_aging_cal,"Balance",2)

priv = template_aging_cal.iloc[row_no,col_no:col_no+9]
ins = template_aging_cal.iloc[row_no+1,col_no:col_no+9]
total = template_aging_cal.iloc[row_no+2,col_no:col_no+9]

for shape in slide.shapes:
    if shape.name == "Chart 4":
        chart = shape.chart
        
chart_data = CategoryChartData()
chart_data.categories = ["Insurance","Private","Total"]
categories = ["0-30","31-60","61-90","91-120","121-150","151-180","181-210","211-365","365+"]

for i in range(len(categories)):
    chart_data.add_series(categories[i], (ins.iloc[i],priv.iloc[i],total.iloc[i]))
chart.replace_data(chart_data)

row_no_start, col_no = coordinate_finder_insheet_df(template_aging_cal,"AR Over 90",1)
row_no_end, col_no_new = coordinate_finder_insheet_df(template_aging_cal,"Payer Type Total",1)

over_90 = template_aging_cal.iloc[row_no_start:row_no_end,col_no-1:col_no+1]
over_90.iloc[:,0] = over_90.iloc[:,0].astype(int)

for i in range(len(over_90)):
    over_90.iloc[i,0] = '${:,}'.format(over_90.iloc[i,0])
    over_90.iloc[i,1] = "{:.1%}".format(over_90.iloc[i,1])
    
ins_type = template_aging_cal.iloc[row_no_start:row_no_end,col_no_new-1]
ar_table = pd.concat([ins_type,over_90],axis=1)

for shape in slide.shapes:
    if shape.name == "Table 4":
        table = shape.table

for i in range(ar_table.shape[0]):
    for j in range(ar_table.shape[1]):
        cell = table.cell(i+1, j)
        cell.text = ar_table.iloc[i,j]
        text_format(cell)
        if i == len(ar_table)-1:
            bold(cell)

rows_to_delete = 10 - len(ar_table)
for i in range(len(ar_table)+rows_to_delete,len(ar_table),-1):
    row = table.rows[i]
    remove_row(table,row)


#--------------------------------------------------------------------Slide 11
slide = root.slides[10]

for shape in slide.shapes:
    if shape.name == "Table 7":
        table = shape.table
   
values = [ncr,ncr_qhcr,ncr_wo_private,(0.98-ncr_wo_private),ar,ar_qhcr,dso,dso_qhcr]

for i in range(1,9):
    cell = table.cell(2,i)
    if i == 7 : 
        cell.text = str(values[i-1])
    elif i == 8 : 
        if values[i-1] > 0 : 
            cell.text = "-"
        else:
            cell.text = str(abs(values[i-1]))
        font_color(cell,"green")
    elif i in range(1,5):
        cell.text = "{:.1%}".format(abs(values[i-1]))
        if i%2 == 0 :
            if values[i-1] < 0 : 
                cell.text = "-"
            font_color(cell,"green")
    else :
        cell.text = "{:.1%}".format(abs(values[i-1]))
        if i%2 == 0 :
            if values[i-1] > 0 : 
                cell.text = "-"
            font_color(cell,"green")
        
    text_format(cell)
    
for i in range(1,9):
    cell = table.cell(1,i)
    if i%2 != 0 :
        cell.text = assessment_name
        text_format_heading(cell)

root.save(output_path)

ppt = win32com.client.Dispatch('PowerPoint.Application')
root = ppt.Presentations.Open(output_path)
ppt.Run("powerpoint.pptm!module1.main")
root.SaveAs(final_path)
root.close()
ppt.Quit()

